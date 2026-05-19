from __future__ import annotations

import json
import math
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "card" / "final" / "v2" / "参赛文件"
ASSET_DIR = ROOT / "output" / "v2_live_selection_ppt"
PREVIEW_DIR = ASSET_DIR / "previews"
ILLUSTRATION_DIR = ROOT / "card" / "final" / "v2" / "generated_illustrations"
PPTX_PATH = OUT_DIR / "《实验室安全值班》v2化学实验室安全情境协作桌游---赵成博.pptx"

WIDE_W = 13.333
WIDE_H = 7.5

BLUE = RGBColor(22, 58, 112)
MID_BLUE = RGBColor(64, 127, 216)
DARK = RGBColor(25, 34, 48)
GRAY = RGBColor(85, 96, 112)
LIGHT = RGBColor(241, 246, 252)
ORANGE = RGBColor(240, 161, 58)
GREEN = RGBColor(47, 164, 143)
PURPLE = RGBColor(126, 91, 216)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    paths = [
        Path(r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\simsun.ttc"),
    ]
    for path in paths:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, line_spacing=1.05):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=18, color=DARK):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_image(slide, path: Path, x, y, w=None, h=None):
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), height=Cm(h))
    return slide.shapes.add_picture(str(path), Cm(x), Cm(y))


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 251, 255)
    return slide


def make_montage(ids: list[str], out_path: Path, cols: int = 6, thumb_w: int = 220, thumb_h: int = 330) -> Path:
    rows = math.ceil(len(ids) / cols)
    label_h = 34
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + label_h)), "#F4F7FB")
    draw = ImageDraw.Draw(sheet)
    font = load_font(18, bold=True)
    for idx, cid in enumerate(ids):
        path = ILLUSTRATION_DIR / f"{cid}.png"
        if not path.exists():
            continue
        img = Image.open(path).convert("RGB")
        img.thumbnail((thumb_w - 12, thumb_h - 10), Image.Resampling.LANCZOS)
        col, row = idx % cols, idx // cols
        x = col * thumb_w + (thumb_w - img.width) // 2
        y = row * (thumb_h + label_h) + 2
        sheet.paste(img, (x, y))
        bbox = draw.textbbox((0, 0), cid, font=font)
        draw.text((col * thumb_w + (thumb_w - (bbox[2] - bbox[0])) / 2, y + thumb_h + 4), cid, fill="#172033", font=font)
    sheet.save(out_path, quality=92)
    return out_path


def make_chart(out_path: Path) -> Path:
    font_path = next((path for path in [Path(r"C:\Windows\Fonts\msyh.ttc"), Path(r"C:\Windows\Fonts\simhei.ttf")] if path.exists()), None)
    if font_path:
        font_manager.fontManager.addfont(str(font_path))
        plt.rcParams["font.sans-serif"] = [font_manager.FontProperties(fname=str(font_path)).get_name()]
    plt.rcParams["axes.unicode_minus"] = False
    data = json.loads((ROOT / "data" / "v2" / "run_009_v2_prefinal_longrun" / "summary.json").read_text(encoding="utf-8"))
    rows = {row["mode"]: row for row in data["scenario_stats"] if row["strategy_id"] == "none"}
    labels = ["教学", "标准", "挑战"]
    values = [rows["teaching"]["win_rate"] * 100, rows["standard"]["win_rate"] * 100, rows["challenge"]["win_rate"] * 100]
    plt.figure(figsize=(8, 4.2), dpi=180)
    bars = plt.bar(labels, values, color=["#3F7FD8", "#2FA48F", "#F0A13A"], width=0.55)
    plt.ylim(0, 100)
    plt.ylabel("胜率 (%)")
    plt.title("v2 长跑模拟：难度梯度清晰")
    for bar, value in zip(bars, values):
        plt.text(bar.get_x() + bar.get_width() / 2, value + 2, f"{value:.1f}%", ha="center", fontsize=12)
    plt.tight_layout()
    plt.savefig(out_path)
    plt.close()
    return out_path


def make_slide_preview(path: Path, title: str, subtitle: str, img_paths: list[Path] | None = None) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    if path.exists():
        path.unlink()
    canvas = Image.new("RGB", (1600, 900), "#F8FBFF")
    draw = ImageDraw.Draw(canvas)
    title_font = load_font(56, bold=True)
    body_font = load_font(28)
    small_font = load_font(22)
    draw.text((90, 70), title, fill="#163A70", font=title_font)
    draw.text((92, 150), subtitle, fill="#44546A", font=body_font)
    if img_paths:
        x = 90
        for idx, img_path in enumerate(img_paths[:4]):
            img = Image.open(img_path).convert("RGB")
            img.thumbnail((310, 465), Image.Resampling.LANCZOS)
            canvas.paste(img, (x, 265))
            draw.text((x + 105, 745), img_path.stem, fill="#172033", font=small_font)
            x += 350
    canvas.save(path)
    return path


def build_deck() -> dict:
    ensure_dirs()
    chart_path = make_chart(ASSET_DIR / "winrate_chart.png")
    montage_events = make_montage([f"E{i:02d}" for i in range(1, 37)], ASSET_DIR / "events_montage.jpg", cols=9, thumb_w=160, thumb_h=240)
    montage_all = make_montage(
        [f"E{i:02d}" for i in range(1, 37)]
        + [f"T{i:02d}" for i in range(1, 19)]
        + [f"A{i:02d}" for i in range(1, 25)]
        + [f"R{i:02d}" for i in range(1, 7)]
        + [f"P{i:02d}" for i in range(1, 5)]
        + [f"S{i:02d}" for i in range(1, 9)]
        + [f"D{i:02d}" for i in range(1, 25)],
        ASSET_DIR / "all_cards_montage.jpg",
        cols=12,
        thumb_w=120,
        thumb_h=180,
    )

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # 1. Cover
    slide = blank_slide(prs)
    add_text(slide, "《实验室安全值班》v2", 1.5, 1.2, 15.5, 1.4, 34, True, BLUE)
    add_text(slide, "化学实验室安全情境协作桌游", 1.55, 2.55, 13.0, 0.8, 18, False, GRAY)
    add_text(slide, "4 分钟现场汇报  |  汇报人：赵成博", 1.55, 3.35, 13.5, 0.7, 15, False, GRAY)
    add_image(slide, ILLUSTRATION_DIR / "E02.png", 19.1, 0.55, h=8.4)
    add_image(slide, ILLUSTRATION_DIR / "E15.png", 24.25, 1.2, h=7.4)
    add_text(slide, "把“知道规范”变成“情境中做对选择”", 1.55, 11.5, 17.0, 1.2, 24, True, DARK)
    add_bullets(slide, ["任务推进", "风险事件", "岗位协作", "知识复盘"], 1.55, 13.0, 13.0, 2.0, 18, BLUE)

    # 2. Problem
    slide = blank_slide(prs)
    add_text(slide, "为什么要做成桌游？", 1.2, 0.9, 18, 1.2, 30, True, BLUE)
    add_text(slide, "实验室安全的难点不是背条文，而是在具体情境中判断优先级。", 1.25, 2.2, 21.0, 0.9, 20, True, DARK)
    add_bullets(slide, ["风险经常是渐进出现：标签不清、通风不足、装置松动。", "真实处置需要协作：谁停止操作、谁记录、谁上报。", "复盘比处罚更重要：知道为什么这样做，下一次才会做对。"], 1.3, 3.6, 16.5, 4.5, 20)
    add_image(slide, ILLUSTRATION_DIR / "E01.png", 21.8, 1.2, h=7.3)
    add_image(slide, ILLUSTRATION_DIR / "E05.png", 26.6, 3.0, h=6.6)

    # 3. What it is
    slide = blank_slide(prs)
    add_text(slide, "作品是什么：一套协作式安全训练系统", 1.2, 0.8, 21.0, 1.0, 28, True, BLUE)
    add_text(slide, "2-4 人  |  20-30 分钟  |  120 张结构化卡牌", 1.25, 2.0, 18.0, 0.8, 20, True, GREEN)
    add_bullets(slide, ["任务牌：推进真实实验流程", "事件牌：触发常见实验室风险", "行动牌：选择防护、核对、通风、停机、隔离、上报", "复盘牌：把处置动作连接到安全知识点"], 1.3, 3.25, 16.0, 5.5, 19)
    add_image(slide, montage_events, 19.0, 1.25, w=13.2)

    # 4. Core loop
    slide = blank_slide(prs)
    add_text(slide, "一轮怎么玩：4 个动作形成教学闭环", 1.2, 0.8, 20.0, 1.0, 28, True, BLUE)
    steps = [("1", "推进任务", "领取试剂、搭建装置、转移溶剂"), ("2", "遭遇事件", "标签、通风、加热、废液、泄漏"), ("3", "岗位处置", "安全员、操作者、记录员协作"), ("4", "复盘知识", "正确处置、原因、反思问题")]
    x = 1.3
    for idx, (num, head, body) in enumerate(steps):
        color = [MID_BLUE, ORANGE, GREEN, PURPLE][idx]
        add_rect(slide, x, 3.0, 6.8, 4.3, RGBColor(234, 242, 255))
        add_text(slide, num, x + 0.35, 3.22, 0.8, 0.8, 26, True, color)
        add_text(slide, head, x + 1.25, 3.23, 4.9, 0.7, 22, True, BLUE)
        add_text(slide, body, x + 0.45, 4.45, 5.8, 1.5, 17, False, DARK)
        if idx < 3:
            add_text(slide, "→", x + 6.95, 4.2, 1.0, 0.8, 30, True, GRAY, PP_ALIGN.CENTER)
        x += 7.55
    add_text(slide, "关键：玩家不是回答“对/错”，而是在有限行动中讨论“先做什么”。", 1.4, 9.0, 26.0, 1.0, 21, True, DARK)

    # 5. Maturity
    slide = blank_slide(prs)
    add_text(slide, "v2 的成熟化：从卡牌原型到中量桌游", 1.2, 0.8, 22.0, 1.0, 28, True, BLUE)
    add_bullets(slide, ["风险状态：苗头 → 暴露 → 失控，不再只是一张牌扣分。", "岗位轮值：安全员、操作者、记录员、资源管理员各有职责。", "交接链：未处理干净的隐患会影响下一轮。", "个人贡献：团队共同胜负，但每个玩家都有参与感。"], 1.3, 2.3, 16.0, 5.5, 20)
    add_image(slide, ILLUSTRATION_DIR / "P01.png", 20.5, 1.05, h=7.4)
    add_image(slide, ILLUSTRATION_DIR / "R02.png", 25.0, 3.4, h=6.2)

    # 6. Data
    slide = blank_slide(prs)
    add_text(slide, "不是凭感觉调规则：270000 局模拟验证", 1.2, 0.8, 22.0, 1.0, 28, True, BLUE)
    add_image(slide, chart_path, 1.4, 2.0, w=14.6)
    add_bullets(slide, ["教学模式：高胜率，适合新手入门。", "标准模式：有压力但可通过协作获胜。", "挑战模式：熟悉后仍有失败风险。", "平均每局触发 6+ 个知识点。"], 18.0, 2.2, 12.8, 5.0, 20)
    add_text(slide, "数据定位：证明“可玩、可解释、难度有梯度”，不替代真人试玩。", 1.5, 9.2, 25.0, 0.8, 18, False, GRAY)

    # 7. Visuals
    slide = blank_slide(prs)
    add_text(slide, "当前美术：120 张无文字插画已生成", 1.2, 0.8, 21.0, 1.0, 28, True, BLUE)
    add_text(slide, "所有中文后期排版，避免生图模型直接生成文字导致糊字。", 1.25, 1.95, 21.0, 0.8, 19, True, DARK)
    add_image(slide, montage_all, 1.2, 3.0, w=30.8)
    add_text(slide, "蓝白水彩实验室风格：适合平板预览，也可继续进入印刷排版。", 1.3, 16.8, 27.0, 0.8, 18, False, GRAY)

    # 8. Teaching use
    slide = blank_slide(prs)
    add_text(slide, "怎么用：新生入组、课堂讨论、课题组培训", 1.2, 0.8, 23.0, 1.0, 28, True, BLUE)
    add_bullets(slide, ["新生入组：用 20-30 分钟熟悉常见风险。", "课堂讨论：把抽象规范变成具体情境判断。", "课题组培训：按真实实验流程替换任务与事件。", "复赛扩展：补真人试玩、打印样张、制作完整展示视频。"], 1.4, 2.4, 16.0, 5.6, 20)
    add_image(slide, ILLUSTRATION_DIR / "A08.png", 20.5, 1.0, h=7.4)
    add_text(slide, "汇报结束，谢谢各位老师。", 1.4, 9.5, 20.0, 1.0, 26, True, BLUE)

    prs.save(PPTX_PATH)

    preview_specs = [
        ("slide01.png", "《实验室安全值班》v2", "化学实验室安全情境协作桌游", [ILLUSTRATION_DIR / "E02.png", ILLUSTRATION_DIR / "E15.png"]),
        ("slide02.png", "为什么要做成桌游？", "把安全规范转化为情境判断", [ILLUSTRATION_DIR / "E01.png", ILLUSTRATION_DIR / "E05.png"]),
        ("slide07.png", "120 张无文字插画已生成", "平板可快速浏览美术与工作量", [ILLUSTRATION_DIR / "E02.png", ILLUSTRATION_DIR / "T03.png", ILLUSTRATION_DIR / "A08.png", ILLUSTRATION_DIR / "P01.png"]),
    ]
    previews = [make_slide_preview(PREVIEW_DIR / name, title, subtitle, imgs) for name, title, subtitle, imgs in preview_specs]
    return {"pptx": str(PPTX_PATH), "previews": [str(p) for p in previews], "assets": str(ASSET_DIR)}


if __name__ == "__main__":
    print(json.dumps(build_deck(), ensure_ascii=False, indent=2))
