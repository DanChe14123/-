from __future__ import annotations

import json
import math
import re
import shutil
from collections import Counter
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "card" / "final" / "v2" / "参赛文件"
ASSET_DIR = ROOT / "output" / "v2_live_selection_ppt_enhanced"
PREVIEW_DIR = ASSET_DIR / "previews"
ILLUSTRATION_DIR = ROOT / "card" / "final" / "v2" / "generated_illustrations"
PPTX_PATH = OUT_DIR / "《实验室安全值班》v2化学实验室安全情境协作桌游---赵成博.pptx"
BACKUP_PATH = OUT_DIR / "《实验室安全值班》v2化学实验室安全情境协作桌游---赵成博_原简版备份.pptx"

W_CM = 33.867
H_CM = 19.05

BLUE = RGBColor(22, 58, 112)
MID_BLUE = RGBColor(65, 126, 213)
LIGHT_BLUE = RGBColor(231, 241, 255)
DARK = RGBColor(25, 34, 48)
GRAY = RGBColor(86, 99, 118)
LIGHT = RGBColor(248, 251, 255)
PANEL = RGBColor(240, 246, 253)
ORANGE = RGBColor(239, 160, 60)
GREEN = RGBColor(47, 164, 143)
PURPLE = RGBColor(123, 92, 211)
RED = RGBColor(205, 89, 82)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)


def backup_existing() -> None:
    if PPTX_PATH.exists() and not BACKUP_PATH.exists():
        shutil.copy2(PPTX_PATH, BACKUP_PATH)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        Path(r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\simsun.ttc"),
    ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=15, color=DARK, bold_first=False, bullet=False):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ("· " + line) if bullet else line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = bool(bold_first and idx == 0)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def add_panel(slide, x, y, w, h, fill=PANEL, line=RGBColor(217, 228, 243)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_kpi(slide, value, label, x, y, w, h, color=BLUE, note=None):
    add_panel(slide, x, y, w, h, RGBColor(247, 250, 255))
    add_text(slide, value, x + 0.25, y + 0.22, w - 0.5, 0.85, 28, True, color, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.25, y + 1.12, w - 0.5, 0.5, 11, False, GRAY, PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + 0.25, y + 1.62, w - 0.5, 0.45, 9, False, GRAY, PP_ALIGN.CENTER)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 1.15, 0.65, 25.5, 0.9, 27, True, BLUE)
    if subtitle:
        add_text(slide, subtitle, 1.18, 1.58, 24.5, 0.6, 14, False, GRAY)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    return slide


def add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        return None
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), height=Cm(h))
    return slide.shapes.add_picture(str(path), Cm(x), Cm(y))


def make_montage(ids: list[str], out_path: Path, cols: int, thumb_w: int, thumb_h: int) -> Path:
    rows = math.ceil(len(ids) / cols)
    label_h = 28
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + label_h)), "#F6FAFF")
    draw = ImageDraw.Draw(sheet)
    font = load_font(16, bold=True)
    for idx, cid in enumerate(ids):
        p = ILLUSTRATION_DIR / f"{cid}.png"
        if not p.exists():
            continue
        img = Image.open(p).convert("RGB")
        img.thumbnail((thumb_w - 12, thumb_h - 10), Image.Resampling.LANCZOS)
        col, row = idx % cols, idx // cols
        x = col * thumb_w + (thumb_w - img.width) // 2
        y = row * (thumb_h + label_h) + 2
        sheet.paste(img, (x, y))
        bbox = draw.textbbox((0, 0), cid, font=font)
        draw.text((col * thumb_w + (thumb_w - (bbox[2] - bbox[0])) / 2, y + thumb_h + 1), cid, fill="#172033", font=font)
    sheet.save(out_path, quality=92)
    return out_path


def make_slide_preview(path: Path, title: str, subtitle: str, img_ids: list[str] | None = None) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    if path.exists():
        path.unlink()
    canvas = Image.new("RGB", (1600, 900), "#F8FBFF")
    draw = ImageDraw.Draw(canvas)
    title_font = load_font(54, bold=True)
    body_font = load_font(28)
    small_font = load_font(22, bold=True)
    draw.text((80, 65), title, fill="#163A70", font=title_font)
    draw.text((82, 145), subtitle, fill="#44546A", font=body_font)
    if img_ids:
        x = 90
        for cid in img_ids[:4]:
            p = ILLUSTRATION_DIR / f"{cid}.png"
            if not p.exists():
                continue
            img = Image.open(p).convert("RGB")
            img.thumbnail((310, 465), Image.Resampling.LANCZOS)
            canvas.paste(img, (x, 265))
            draw.text((x + 120, 750), cid, fill="#172033", font=small_font)
            x += 350
    canvas.save(path)
    return path


def load_data() -> dict:
    cards = json.loads((ROOT / "content" / "v2_cards.json").read_text(encoding="utf-8"))["cards"]
    counts = Counter(c["card_type"] for c in cards)
    v2 = json.loads((ROOT / "data" / "v2" / "run_009_v2_prefinal_longrun" / "summary.json").read_text(encoding="utf-8"))
    rows = {row["mode"]: row for row in v2["scenario_stats"] if row["strategy_id"] == "none"}
    strategy_standard = [row for row in v2["scenario_stats"] if row["mode"] == "standard"]
    top_actions = sorted(v2["action_pick_rates"].items(), key=lambda kv: kv[1], reverse=True)[:6]
    doc_files = [
        "boardgame_v2.0.md",
        "docs/v2_rule_design_notes.md",
        "content/v2_card_schema.md",
        "content/v2_cards.json",
        "manual/player_rulebook_v2.md",
        "manual/teaching_guide_v2.md",
        "manual/judge_brief_v2.md",
        "manual/quick_reference_v2.md",
        "docs/v2_ai_agent_playtest_protocol.md",
        "docs/v2_ai_agent_playtest_findings.md",
        "docs/v2_sample_playthrough_standard.md",
        "docs/v2_120_card_image_generation_guide.md",
    ]
    doc_chars = sum(len((ROOT / f).read_text(encoding="utf-8", errors="ignore")) for f in doc_files if (ROOT / f).exists())
    v2_total = 0
    for p in (ROOT / "data" / "v2").glob("run_*/summary.json"):
        data = json.loads(p.read_text(encoding="utf-8"))
        v2_total += int(data.get("games_per_scenario", 0)) * len(data.get("scenario_stats", []))
    v0_total = 0
    for p in (ROOT / "data").glob("v0.*/*/summary.json"):
        data = json.loads(p.read_text(encoding="utf-8"))
        v0_total += int(data.get("games_per_team", 0)) * int(data.get("scenarios_tested", 0))
    version_points = [
        ("v0.1", 11.17),
        ("v0.2", 56.27),
        ("v0.5", 51.33),
        ("v0.8", 52.44),
        ("v2", rows["standard"]["win_rate"] * 100),
    ]
    return {
        "counts": counts,
        "rows": rows,
        "strategy_standard": strategy_standard,
        "top_actions": top_actions,
        "doc_chars": doc_chars,
        "v2_total": v2_total,
        "v0_total": v0_total,
        "version_points": version_points,
    }


def style_chart(chart, font_size=10):
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(font_size)
    chart.category_axis.tick_labels.font.name = "Microsoft YaHei"
    chart.value_axis.tick_labels.font.size = Pt(font_size)
    chart.value_axis.tick_labels.font.name = "Microsoft YaHei"
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(220, 230, 242)


def add_column_chart(slide, categories, values, x, y, w, h, color=BLUE, max_scale=None):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("胜率", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), data).chart
    style_chart(chart)
    chart.value_axis.maximum_scale = max_scale if max_scale is not None else max(values) * 1.25
    chart.value_axis.minimum_scale = 0
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.has_data_labels = True
    series.data_labels.number_format = '0.0'
    series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return chart


def add_line_chart(slide, categories, values, x, y, w, h):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("标准模式胜率", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Cm(x), Cm(y), Cm(w), Cm(h), data).chart
    style_chart(chart)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 75
    series = chart.series[0]
    series.format.line.color.rgb = MID_BLUE
    series.format.line.width = Pt(2.5)
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = ORANGE
    series.has_data_labels = True
    series.data_labels.number_format = '0.0'
    series.data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
    return chart


def add_table(slide, rows, x, y, w, h, col_widths=None, header_fill=BLUE):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Cm(x), Cm(y), Cm(w), Cm(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Cm(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.04)
            cell.margin_bottom = Cm(0.04)
            fill = header_fill if r == 0 else (RGBColor(247, 250, 255) if r % 2 else RGBColor(238, 245, 253))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(10 if len(str(value)) < 22 else 8.5)
                    run.font.bold = r == 0
                    run.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else DARK
    return table_shape


def build_deck() -> dict:
    ensure_dirs()
    backup_existing()
    data = load_data()

    event_montage = make_montage([f"E{i:02d}" for i in range(1, 37)], ASSET_DIR / "event_montage.jpg", 9, 150, 225)
    all_montage = make_montage(
        [f"E{i:02d}" for i in range(1, 37)]
        + [f"T{i:02d}" for i in range(1, 19)]
        + [f"A{i:02d}" for i in range(1, 25)]
        + [f"R{i:02d}" for i in range(1, 7)]
        + [f"P{i:02d}" for i in range(1, 5)]
        + [f"S{i:02d}" for i in range(1, 9)]
        + [f"D{i:02d}" for i in range(1, 25)],
        ASSET_DIR / "all_montage.jpg",
        12,
        110,
        165,
    )

    prs = Presentation()
    prs.slide_width = Cm(W_CM)
    prs.slide_height = Cm(H_CM)

    # 1 Cover
    slide = blank_slide(prs)
    add_text(slide, "《实验室安全值班》v2", 1.1, 0.9, 15.5, 1.2, 34, True, BLUE)
    add_text(slide, "化学实验室安全情境协作桌游", 1.15, 2.08, 12.5, 0.6, 17, False, GRAY)
    add_text(slide, "现场评选 4 分钟汇报  |  汇报人：赵成博", 1.15, 2.82, 13.2, 0.55, 13, False, GRAY)
    add_image(slide, ILLUSTRATION_DIR / "E02.png", 19.3, 0.55, h=8.3)
    add_image(slide, ILLUSTRATION_DIR / "E15.png", 24.3, 1.15, h=7.35)
    add_text(slide, "把“知道规范”变成“情境中做对选择”", 1.15, 10.15, 15.5, 0.9, 24, True, DARK)
    add_kpi(slide, "120", "结构化卡牌", 1.2, 11.65, 4.3, 2.0, MID_BLUE)
    add_kpi(slide, "27万", "最终长跑模拟", 5.9, 11.65, 4.3, 2.0, GREEN)
    add_kpi(slide, "8+1", "v0.1-v0.8 + v2", 10.6, 11.65, 4.3, 2.0, ORANGE)
    add_kpi(slide, "35万+", "规则/文案/提示词字符", 15.3, 11.65, 5.0, 2.0, PURPLE)

    # 2 Problem and answer
    slide = blank_slide(prs)
    add_title(slide, "为什么要做成桌游？", "安全培训的难点不是背条文，而是在具体情境中判断优先级")
    add_image(slide, ILLUSTRATION_DIR / "E05.png", 1.15, 2.6, h=7.0)
    add_image(slide, ILLUSTRATION_DIR / "A08.png", 6.15, 3.2, h=6.35)
    add_panel(slide, 13.3, 2.55, 8.6, 6.6)
    add_text(slide, "真实实验中的决策压力", 13.75, 2.95, 7.7, 0.5, 17, True, BLUE)
    add_multiline(slide, [
        "风险经常是渐进出现：标签不清、通风不足、装置松动。",
        "处置不是单人答题：需要停止、隔离、记录、上报和复盘。",
        "同一个事件有先后顺序：先控风险，再推进任务。",
    ], 13.75, 3.85, 7.6, 2.5, 13.5, DARK, bullet=True)
    add_panel(slide, 22.8, 2.55, 8.9, 6.6, RGBColor(237, 248, 246))
    add_text(slide, "桌游带来的训练价值", 23.25, 2.95, 8.0, 0.5, 17, True, GREEN)
    add_multiline(slide, [
        "把抽象规范压缩成一轮一轮的选择。",
        "让玩家在低风险环境中练习协作判断。",
        "每局结束都有知识点复盘，形成教学闭环。",
    ], 23.25, 3.85, 7.8, 2.5, 13.5, DARK, bullet=True)

    # 3 Iteration timeline
    slide = blank_slide(prs)
    add_title(slide, "版本不是一次做完：从 v0.1 到 v2 的迭代轨迹", "规则、牌池、模拟器、说明书和美术层逐步拆分")
    timeline = [
        ("v0.1", "最小可玩原型", "验证任务推进 + 事件处置"),
        ("v0.3", "加入隐患压力", "风险不再只是扣分"),
        ("v0.5", "任务牌扩充", "实验流程从零散事件变成路线"),
        ("v0.8", "初赛封版", "56 张卡 + 长跑数据 + 提交包"),
        ("v2", "成熟化重构", "120 张卡 + 岗位轮值 + 交接链"),
    ]
    x = 1.2
    for idx, (ver, head, body) in enumerate(timeline):
        color = [BLUE, MID_BLUE, GREEN, ORANGE, PURPLE][idx]
        add_panel(slide, x, 3.0, 5.75, 5.1, RGBColor(247, 250, 255))
        add_text(slide, ver, x + 0.35, 3.35, 4.9, 0.6, 22, True, color, PP_ALIGN.CENTER)
        add_text(slide, head, x + 0.35, 4.35, 4.9, 0.55, 15, True, DARK, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.55, 5.35, 4.5, 1.25, 12.5, False, GRAY, PP_ALIGN.CENTER)
        if idx < len(timeline) - 1:
            add_text(slide, "→", x + 5.82, 5.0, 0.8, 0.8, 24, True, GRAY, PP_ALIGN.CENTER)
        x += 6.35
    add_line_chart(slide, [p[0] for p in data["version_points"]], [p[1] for p in data["version_points"]], 2.2, 10.0, 18.0, 5.4)
    add_text(slide, "标准模式胜率从“过难”收敛到可玩区间；v2 不是单纯加牌，而是重构了协作结构。", 21.0, 10.15, 10.6, 2.0, 15, True, DARK)
    add_text(slide, "注：v2 规则体系已重构，曲线用于展示难度收敛趋势，不作为严格同模型对比。", 21.0, 12.15, 10.6, 0.85, 10, False, GRAY)

    # 4 Rule innovations
    slide = blank_slide(prs)
    add_title(slide, "玩法成熟化：v2 的 5 个核心创新点", "从“抽事件回答”升级为“实验任务层 + 风险处置层 + 岗位协作层”")
    innovations = [
        ("风险状态", "苗头 → 暴露 → 失控，让事件有发展过程。"),
        ("岗位轮值", "安全员、操作者、记录员、资源管理员各有职责。"),
        ("交接链", "上一轮没处理干净的隐患会影响下一轮。"),
        ("贡献分", "团队共同胜负，但每名玩家有参与反馈。"),
        ("复盘牌", "把正确处置、原因和知识点独立沉淀。"),
    ]
    for i, (head, body) in enumerate(innovations):
        y = 2.65 + i * 2.35
        add_text(slide, f"{i + 1}", 1.35, y, 0.8, 0.6, 18, True, [MID_BLUE, GREEN, ORANGE, PURPLE, RED][i], PP_ALIGN.CENTER)
        add_text(slide, head, 2.25, y - 0.02, 4.2, 0.55, 18, True, BLUE)
        add_text(slide, body, 6.2, y, 11.5, 0.55, 13.5, False, DARK)
    add_image(slide, ILLUSTRATION_DIR / "P01.png", 20.0, 2.05, h=7.15)
    add_image(slide, ILLUSTRATION_DIR / "R02.png", 25.15, 4.05, h=6.25)
    add_text(slide, "设计目标：让评委看到它不是“安全知识卡片”，而是一套有回合、资源、协作与复盘结构的桌游。", 19.9, 11.35, 11.4, 1.2, 14.5, True, DARK)

    # 5 Card system
    slide = blank_slide(prs)
    add_title(slide, "最终体量：120 张卡牌，覆盖完整教学闭环", "卡牌不是堆数量，而是按功能分工支撑一局游戏")
    count_labels = ["事件", "行动", "复盘", "任务", "策略", "角色", "岗位"]
    count_values = [
        data["counts"]["event"],
        data["counts"]["action"],
        data["counts"]["debrief"],
        data["counts"]["task"],
        data["counts"]["strategy"],
        data["counts"]["role"],
        data["counts"]["post"],
    ]
    add_column_chart(slide, count_labels, count_values, 1.4, 2.65, 14.2, 6.0, MID_BLUE, max_scale=40)
    rows = [
        ["类型", "数量", "作用"],
        ["任务牌", 18, "组织实验流程与胜利进度"],
        ["事件牌", 36, "触发常见实验室风险"],
        ["行动牌", 24, "提供处置选择"],
        ["复盘牌", 24, "承载知识点与提问"],
        ["角色/岗位/策略", 18, "增强协作与打法差异"],
    ]
    add_table(slide, rows, 17.0, 2.65, 14.5, 6.0, [3.2, 2.0, 8.8])
    add_text(slide, "卡牌结构对应“任务推进 → 风险响应 → 正确处置 → 知识复盘”，方便老师解释，也方便学生上手。", 1.45, 10.25, 29.5, 0.9, 16, True, DARK)
    add_image(slide, event_montage, 1.4, 12.0, w=30.0)

    # 6 Writing and material workload
    slide = blank_slide(prs)
    add_title(slide, "文案与材料工作量：不是只有卡图", "规则、文案、说明书、评委入口和 AI 代理测试都有独立文件")
    add_kpi(slide, "35.1万", "v2 规则/文案/提示词字符", 1.25, 2.8, 6.0, 2.3, PURPLE)
    add_kpi(slide, "4 本", "玩家/教师/评委/速查材料", 8.05, 2.8, 6.0, 2.3, GREEN)
    add_kpi(slide, "120 条", "逐卡无文字插画提示词", 14.85, 2.8, 6.0, 2.3, MID_BLUE)
    add_kpi(slide, "4 类", "AI 代理试玩视角", 21.65, 2.8, 6.0, 2.3, ORANGE)
    rows = [
        ["材料层", "已经完成的内容"],
        ["规则层", "boardgame_v2.0 + 规则设计说明"],
        ["牌池层", "v2_cards.json + 卡牌字段结构"],
        ["教学层", "玩家规则书、教师指南、评委摘要、一页速查"],
        ["验证层", "平衡报告、示例战报、AI 代理试玩记录"],
        ["美术层", "120 张插画提示词 + 120 张无文字插画"],
    ]
    add_table(slide, rows, 1.35, 6.3, 17.1, 8.1, [4.0, 12.6], header_fill=PURPLE)
    add_panel(slide, 20.2, 6.3, 10.7, 8.1, RGBColor(245, 248, 255))
    add_text(slide, "文案迭代路径", 20.65, 6.8, 9.8, 0.6, 18, True, BLUE)
    add_multiline(slide, [
        "v0.8：完成 56 张卡与初赛文字说明。",
        "v2：重写规则、牌池和说明书体系。",
        "AI 代理测试后：修订速查卡与玩家规则书。",
        "图像阶段：为 120 张卡逐张写无文字插画指令。",
    ], 20.65, 7.85, 9.6, 4.4, 13.2, DARK, bullet=True)

    # 7 Simulation evidence
    slide = blank_slide(prs)
    add_title(slide, "数据不是装饰：用模拟测试胜率、难度和行动分布", "最终长跑基线：run_009_v2_prefinal_longrun，seed=20260428")
    add_kpi(slide, "27万局", "v2 最终长跑", 1.25, 2.55, 4.7, 2.15, MID_BLUE)
    add_kpi(slide, "45.9万局", "v2 已保存模拟记录", 6.45, 2.55, 5.4, 2.15, GREEN)
    add_kpi(slide, "132万+局", "v0.x + v2 累计记录", 12.35, 2.55, 5.4, 2.15, ORANGE)
    add_kpi(slide, "6.51", "标准模式平均知识点/局", 18.25, 2.55, 5.4, 2.15, PURPLE)
    modes = ["教学", "标准", "挑战"]
    values = [
        data["rows"]["teaching"]["win_rate"] * 100,
        data["rows"]["standard"]["win_rate"] * 100,
        data["rows"]["challenge"]["win_rate"] * 100,
    ]
    add_column_chart(slide, modes, values, 1.45, 6.1, 12.3, 7.0, GREEN, max_scale=100)
    rows = [
        ["指标", "结果", "解读"],
        ["教学模式胜率", "82.86%", "适合新手入门"],
        ["标准模式胜率", "59.12%", "有压力但可获胜"],
        ["挑战模式胜率", "31.89%", "熟悉后仍有挑战"],
        ["知识点/局", "6.51", "每局能触发复盘内容"],
    ]
    add_table(slide, rows, 15.2, 6.1, 15.3, 6.5, [4.0, 3.0, 7.8], header_fill=GREEN)
    add_text(slide, "说明：模拟用于排查过难、过易和万能行动牌，不替代后续真人试玩。", 15.25, 13.1, 15.1, 0.7, 12, False, GRAY)

    # 8 Balance details
    slide = blank_slide(prs)
    add_title(slide, "平衡性检查：没有单一行动垄断解法", "行动选取率、策略胜率和岗位触发都被纳入检查")
    action_labels = [a for a, _ in data["top_actions"]]
    action_values = [v * 100 for _, v in data["top_actions"]]
    add_column_chart(slide, action_labels, action_values, 1.4, 2.5, 14.5, 6.2, MID_BLUE, max_scale=18)
    add_text(slide, "最高行动选取率仅 6.20%，低于 18% 阈值，说明没有明显“万能牌”。", 1.55, 9.1, 14.0, 0.8, 14.2, True, DARK)
    std_rates = [r["win_rate"] * 100 for r in data["strategy_standard"]]
    min_rate, max_rate = min(std_rates), max(std_rates)
    add_panel(slide, 17.25, 2.45, 13.3, 3.0, RGBColor(245, 248, 255))
    add_text(slide, "策略胜率范围", 17.65, 2.85, 5.0, 0.5, 16, True, BLUE)
    add_text(slide, f"{min_rate:.2f}% - {max_rate:.2f}%", 17.65, 3.55, 8.6, 0.8, 26, True, ORANGE)
    add_text(slide, "S01 稳健管理略强，S02 效率优先更激进；这是复赛可继续精修点。", 17.65, 4.35, 11.7, 0.65, 11.5, False, GRAY)
    post_total = 493760 + 668419 + 595304 + 302516
    post_rows = [
        ["岗位", "触发占比"],
        ["记录员", f"{668419 / post_total * 100:.1f}%"],
        ["安全员", f"{595304 / post_total * 100:.1f}%"],
        ["资源管理员", f"{493760 / post_total * 100:.1f}%"],
        ["操作者", f"{302516 / post_total * 100:.1f}%"],
    ]
    add_table(slide, post_rows, 17.25, 6.3, 9.7, 5.8, [5.0, 4.2], header_fill=MID_BLUE)
    add_text(slide, "平衡逻辑：如果某个行动、岗位或策略过强，模拟数据会先暴露，再进入下一轮调参。", 17.25, 12.7, 12.5, 1.0, 13.5, True, DARK)

    # 9 Visual production
    slide = blank_slide(prs)
    add_title(slide, "美术生产：120 张无文字插画层已经生成", "中文后期排版，避免生图模型直接生成文字导致糊字")
    add_image(slide, all_montage, 1.2, 2.55, w=19.5)
    add_panel(slide, 22.1, 2.55, 8.9, 10.8, RGBColor(245, 248, 255))
    add_text(slide, "统一视觉标准", 22.55, 3.0, 7.8, 0.6, 18, True, BLUE)
    add_multiline(slide, [
        "蓝白色调、现代化实验室、水彩纸纹。",
        "所有插画不含可读文字，后期统一排中文。",
        "角色、场景、器皿风格统一，方便继续印刷化。",
        "当前提交包已含 120 张卡牌插画和总览图。",
    ], 22.55, 4.0, 7.8, 4.6, 13.2, DARK, bullet=True)
    add_kpi(slide, "120", "卡牌插画", 22.55, 9.2, 3.7, 2.1, MID_BLUE)
    add_kpi(slide, "0", "生图中文字", 26.8, 9.2, 3.7, 2.1, GREEN)

    # 10 Landing
    slide = blank_slide(prs)
    add_title(slide, "现场想让老师记住的 3 件事", "作品不是单张海报，而是一套可继续试玩、迭代和教学使用的安全训练桌游")
    add_panel(slide, 1.3, 2.85, 9.2, 7.6, RGBColor(245, 248, 255))
    add_text(slide, "1 主题明确", 1.85, 3.35, 7.9, 0.7, 22, True, BLUE)
    add_text(slide, "围绕化学实验室安全，不泛化成普通科普桌游。", 1.85, 4.45, 7.3, 1.3, 16, False, DARK)
    add_panel(slide, 12.2, 2.85, 9.2, 7.6, RGBColor(241, 249, 247))
    add_text(slide, "2 机制完整", 12.75, 3.35, 7.9, 0.7, 22, True, GREEN)
    add_text(slide, "有任务、事件、行动、岗位、策略、复盘和胜负条件。", 12.75, 4.45, 7.3, 1.3, 16, False, DARK)
    add_panel(slide, 23.1, 2.85, 9.2, 7.6, RGBColor(255, 248, 238))
    add_text(slide, "3 有证据链", 23.65, 3.35, 7.9, 0.7, 22, True, ORANGE)
    add_text(slide, "版本归档、模拟数据、AI 代理测试、完整提交包均可追溯。", 23.65, 4.45, 7.3, 1.3, 16, False, DARK)
    add_text(slide, "适用场景：新生入组培训、课题组安全讨论、实验课前风险演练。", 2.0, 12.1, 27.5, 0.9, 20, True, DARK, PP_ALIGN.CENTER)
    add_text(slide, "谢谢各位老师。", 2.0, 14.2, 27.5, 0.7, 22, True, BLUE, PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)

    previews = [
        make_slide_preview(PREVIEW_DIR / "slide01_cover.png", "《实验室安全值班》v2", "封面加入 120 / 27万 / 8+1 / 35万+ 工作量指标", ["E02", "E15"]),
        make_slide_preview(PREVIEW_DIR / "slide03_iteration.png", "版本迭代轨迹", "v0.1-v0.8 到 v2 的机制与数据收敛", ["E01", "E05", "T03", "A08"]),
        make_slide_preview(PREVIEW_DIR / "slide07_simulation.png", "模拟验证", "最终长跑 27 万局，累计保存模拟记录 132 万+局", ["T03", "E15", "A08", "D02"]),
        make_slide_preview(PREVIEW_DIR / "slide09_visuals.png", "120 张无文字插画", "平板可快速浏览美术与工作量", ["E02", "T03", "A08", "P01"]),
    ]
    return {"pptx": str(PPTX_PATH), "backup": str(BACKUP_PATH), "previews": [str(p) for p in previews]}


if __name__ == "__main__":
    print(json.dumps(build_deck(), ensure_ascii=False, indent=2))
